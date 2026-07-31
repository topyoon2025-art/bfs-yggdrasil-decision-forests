#!/usr/bin/env python3
"""Generate the Symmetric vs Hungarian scaling analysis presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Palette ---
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
BG_SLIDE  = RGBColor(0x16, 0x21, 0x3E)
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2   = RGBColor(0x7B, 0x68, 0xEE)
GREEN     = RGBColor(0x00, 0xE6, 0x96)
RED       = RGBColor(0xFF, 0x6B, 0x6B)
ORANGE    = RGBColor(0xFF, 0xA5, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xB0, 0xB0, 0xC0)
LIGHT     = RGBColor(0xE0, 0xE0, 0xF0)
CARD_BG   = RGBColor(0x1E, 0x2A, 0x4A)
TABLE_HDR = RGBColor(0x0A, 0x14, 0x2A)
TABLE_ROW = RGBColor(0x14, 0x1E, 0x3A)
TABLE_ALT = RGBColor(0x1A, 0x28, 0x48)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width, height=Pt(4), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

add_accent_bar(slide, Inches(0), Inches(0), W, Pt(6), ACCENT)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Why Symmetric Scales and Hungarian Doesn't",
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "GPU Column-Grouping Strategies for Oblique Apply-Projection",
             font_size=24, color=ACCENT)

add_accent_bar(slide, Inches(1), Inches(4.3), Inches(3), Pt(3), ACCENT2)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(1.2),
             "Between 2M×16 and 8M×16 datasets, symmetric’s GPU kernel "
             "performance is stable — but hungarian degrades. Why?",
             font_size=18, color=GRAY)


# ============================================================
# SLIDE 2 — What Each Strategy Does
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             "What Each Strategy Does", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(2.5), Pt(3), ACCENT)

# Left card — Symmetric
add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.5),
             "SYMMETRIC  (strategy 6)", font_size=22, color=GREEN, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5), Inches(3.8), [
    "•  Replace every node’s column indices with a single shared set",
    "•  Pick the K most popular columns across all nodes",
    "•  Replicate that identical assignment to every row",
    "•  Result: 100% column uniformity",
    "•  Cost: O(N·K) — independent of dataset row count",
], font_size=15, color=LIGHT)

# Right card — Hungarian
add_card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
             "HUNGARIAN  (strategy 4)", font_size=22, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(3.8), [
    "•  Reorder each node’s column indices to maximize alignment",
    "•  Per iter: remove node → build K×K cost matrix → solve O(K³) Hungarian → reinsert",
    "•  30 iterations × N nodes per iteration",
    "•  Result: partial column uniformity (best-effort)",
    "•  Cost: O(iters × N × K³) — also independent of row count",
], font_size=15, color=LIGHT)

add_text_box(slide, Inches(0.6), Inches(6.85), Inches(12), Inches(0.5),
             "Neither strategy’s reordering cost depends on #rows. "
             "The difference is what happens in the GPU kernel AFTER reordering.",
             font_size=16, color=GRAY)


# ============================================================
# SLIDE 3 — The GPU Kernel: It's Always a Gather
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             "The GPU Kernel: It’s Always a Gather",
             font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(2.5), Pt(3), ACCENT)

# Code card
add_card(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.0),
         RGBColor(0x0A, 0x12, 0x28))
code_lines = (
    "ex_idx     = selected_examples[row_start + r];          // monotonic, NOT contiguous\n"
    "dataset_idx = col * num_rows + ex_idx;                   // column-major layout\n"
    "x          = dataset[dataset_idx];                       // GATHER — sparse read"
)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(11.4), Inches(1.6),
             code_lines, font_size=16, color=ACCENT, font_name="Consolas")

# Explanation bullets
add_bullet_list(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.5), [
    "•  selected_examples per node are sorted (monotonic) but have gaps between indices",
    "•  Consecutive threads in a warp hit nearby but NOT adjacent memory addresses",
    "•  Each 128-byte cache line (32 floats) may carry only a few useful values",
    "",
    "•  This is true for BOTH symmetric and hungarian — neither gives contiguous reads",
], font_size=17, color=LIGHT)

# Key insight box
add_card(slide, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0),
         RGBColor(0x0A, 0x20, 0x10))
add_text_box(slide, Inches(1.0), Inches(6.05), Inches(11.4), Inches(0.7),
             "Key insight: the advantage of symmetric is NOT \"contiguous reads.\" "
             "It’s about how many DISTINCT columns the L2 cache must hold simultaneously.",
             font_size=18, color=GREEN, bold=True)


# ============================================================
# SLIDE 4 — Working Set: 1 Column vs C Columns
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             "L2 Working Set: 1 Column vs C Columns",
             font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(2.5), Pt(3), ACCENT)

# Left — Symmetric diagram
add_card(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5))
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(5), Inches(0.4),
             "SYMMETRIC — Projection slot j", font_size=20, color=GREEN, bold=True)

sym_lines = (
    "Node  0  ── gather ──▶  column 5\n"
    "Node  1  ── gather ──▶  column 5\n"
    "Node  2  ── gather ──▶  column 5\n"
    "  ...            ...         \n"
    "Node 63  ── gather ──▶  column 5"
)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(2.5),
             sym_lines, font_size=16, color=LIGHT, font_name="Consolas")

add_bullet_list(slide, Inches(1.0), Inches(4.6), Inches(5.2), Inches(2.0), [
    "•  All blocks read SAME column",
    "•  Cache lines shared across blocks",
    "•  L2 footprint = 1 × num_rows × 4B",
], font_size=15, color=LIGHT)

# Right — Hungarian diagram
add_card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.5))
add_text_box(slide, Inches(7.3), Inches(1.55), Inches(5), Inches(0.4),
             "HUNGARIAN — Projection slot j", font_size=20, color=ORANGE, bold=True)

hun_lines = (
    "Node  0  ── gather ──▶  column 5   │ aligned\n"
    "Node  1  ── gather ──▶  column 5   │\n"
    "Node  2  ── gather ──▶  column 12  │ mis-\n"
    "Node  3  ── gather ──▶  column 3   │ aligned\n"
    "  ...            ...         "
)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5.4), Inches(2.5),
             hun_lines, font_size=16, color=LIGHT, font_name="Consolas")

add_bullet_list(slide, Inches(7.3), Inches(4.6), Inches(5.2), Inches(2.0), [
    "•  C distinct columns active at once",
    "•  No cross-node L2 reuse between cols",
    "•  L2 footprint = C × num_rows × 4B",
], font_size=15, color=LIGHT)


# ============================================================
# SLIDE 5 — The Scaling Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             "Why 8M Breaks Hungarian but Not Symmetric",
             font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(2.5), Pt(3), ACCENT)

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "GPU L2 cache: ~40 MB (A100)  /  ~50 MB (H100)",
             font_size=18, color=GRAY)

# Build table manually with cards
table_top = Inches(2.0)
col_widths = [Inches(2.2), Inches(2.2), Inches(3.5), Inches(4.2)]
col_lefts = [Inches(0.6)]
for w in col_widths[:-1]:
    col_lefts.append(col_lefts[-1] + w + Inches(0.05))

row_h = Inches(0.65)
headers = ["Dataset", "1 Column", "Symmetric\n(1 col active)", "Hungarian\n(C≈3 cols active)"]

# Header row
for i, (left, width, text) in enumerate(zip(col_lefts, col_widths, headers)):
    add_card(slide, left, table_top, width, row_h, TABLE_HDR)
    add_text_box(slide, left, table_top, width, row_h,
                 text, font_size=15, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)

# Data rows
data = [
    ("2M × 16", "8 MB", "8 MB   ✓  fits L2", "24 MB   ✓  fits L2"),
    ("8M × 16", "32 MB", "32 MB   ✓  fits L2", "96 MB   ✗  OVERFLOW"),
]
colors_last_col = [GREEN, RED]

for r, (row_data, last_color) in enumerate(zip(data, colors_last_col)):
    y = table_top + (r + 1) * (row_h + Inches(0.05))
    bg = TABLE_ROW if r % 2 == 0 else TABLE_ALT
    for i, (left, width, text) in enumerate(zip(col_lefts, col_widths, row_data)):
        add_card(slide, left, y, width, row_h, bg)
        c = last_color if i == 3 else LIGHT
        add_text_box(slide, left, y, width, row_h,
                     text, font_size=16, color=c, alignment=PP_ALIGN.CENTER)

# Consequences
add_card(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.8),
         RGBColor(0x1A, 0x10, 0x10))
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(11), Inches(0.5),
             "Once L2 overflows:", font_size=20, color=RED, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(2.0), [
    "•  Every cache miss on a gather goes to HBM (~10× slower than L2)",
    "•  Sparse access amplifies the penalty: each evicted 128-byte cache line carried mostly unused bytes",
    "•  Thrashing: column A’s lines evict column B’s, and vice versa — no one stays warm",
], font_size=17, color=LIGHT)


# ============================================================
# SLIDE 6 — Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             "Conclusion", font_size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.6), Inches(1.05), Inches(2.5), Pt(3), ACCENT)

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "•  Both strategies produce a sparse gather — neither gives contiguous reads",
    "",
    "•  The critical difference is L2 working set size:",
    "      Symmetric  →  1 column per projection slot  →  scales linearly with row count",
    "      Hungarian  →  C columns per projection slot  →  scales with C × row count",
    "",
    "•  At 2M rows: even C > 1 columns fit in L2 → both strategies perform similarly",
    "",
    "•  At 8M rows: C > 1 columns overflow L2 → hungarian pays HBM latency",
    "    on every gather miss, while symmetric still fits",
    "",
    "•  Symmetric is dataset-size-insensitive because perfect column uniformity",
    "    confines the working set to one column’s address range, regardless of size",
], font_size=18, color=LIGHT, spacing=Pt(4))


# Save
out_path = "/home/ubuntu/projects/yggdrasil-decision-forests/symmetric_vs_hungarian.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
